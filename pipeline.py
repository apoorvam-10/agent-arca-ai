from typing import List, Dict, Any
import re
import os
import time
import tempfile
from io import BytesIO
from urllib.parse import urlparse

import fitz
import requests
import streamlit as st
from bs4 import BeautifulSoup
import google.generativeai as genai
from tavily import TavilyClient
from youtube_transcript_api import YouTubeTranscriptApi
import webvtt
import yt_dlp

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt as PPTPt

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet


TRUSTED_DOMAINS = [
    "nih.gov",
    "ncbi.nlm.nih.gov",
    "pmc.ncbi.nlm.nih.gov",
    "cdc.gov",
    "who.int",
    "fda.gov",
    "nature.com",
    "sciencedirect.com",
    "springer.com",
    "ieee.org",
    "acm.org",
    "arxiv.org",
    ".edu",
    ".gov",
]

MEDIUM_TRUST_DOMAINS = [
    "ibm.com",
    "microsoft.com",
    "google.com",
    "aws.amazon.com",
    "cloud.google.com",
    "openai.com",
    "wikipedia.org",
    "youtube.com",
    "youtu.be",
]


def get_limits(research_depth: str) -> Dict[str, int]:
    if research_depth == "Deep Research":
        return {
            "web_chars": 10000,
            "url_chars": 10000,
            "video_chars": 12000,
            "file_chars": 12000,
            "image_chars": 10000,
        }

    return {
        "web_chars": 5000,
        "url_chars": 5000,
        "video_chars": 8000,
        "file_chars": 9000,
        "image_chars": 7000,
    }


def get_gemini_model():
    api_key = st.secrets.get("GEMINI_API_KEY", "")

    if not api_key:
        raise ValueError("Missing GEMINI_API_KEY in Streamlit secrets.")

    genai.configure(api_key=api_key)

    return genai.GenerativeModel("gemini-3-flash-preview")


def safe_generate_content(prompt: str) -> str:
    try:
        model = get_gemini_model()
        response = model.generate_content(prompt)

        text = getattr(response, "text", "").strip()

        if not text:
            return "I could not generate a response from the available information."

        return text

    except Exception as e:
        error_text = str(e)

        if (
            "429" in error_text
            or "ResourceExhausted" in error_text
            or "quota" in error_text.lower()
        ):
            return (
                "⚠️ Gemini API quota limit reached. Please wait 1–2 minutes and try again. "
                "This usually happens on the free tier after several requests."
            )

        if "403" in error_text or "leaked" in error_text.lower():
            return (
                "⚠️ Gemini API key issue detected. Your API key may be invalid, blocked, or exposed. "
                "Please create a new key and update Streamlit secrets."
            )

        return f"⚠️ AI generation failed. Please try again. Technical detail: {error_text[:300]}"


def safe_generate_multimodal_content(parts: list) -> str:
    try:
        model = get_gemini_model()
        response = model.generate_content(parts)

        text = getattr(response, "text", "").strip()

        if not text:
            return "I could not extract useful information from this image."

        return text

    except Exception as e:
        error_text = str(e)

        if (
            "429" in error_text
            or "ResourceExhausted" in error_text
            or "quota" in error_text.lower()
        ):
            return "Gemini quota reached while reading image. Please wait 1–2 minutes and try again."

        if "403" in error_text or "leaked" in error_text.lower():
            return "Gemini API key issue detected while reading image. Please update your API key."

        return f"Image reading failed. Technical detail: {error_text[:300]}"


def get_tavily_client():
    api_key = st.secrets.get("TAVILY_API_KEY", "")

    if not api_key:
        raise ValueError("Missing TAVILY_API_KEY in Streamlit secrets.")

    return TavilyClient(api_key=api_key)


def get_domain(url: str) -> str:
    try:
        if not url or not url.startswith("http"):
            return "uploaded-file"

        parsed = urlparse(url)

        return parsed.netloc.lower().replace("www.", "")

    except Exception:
        return "unknown"


def get_source_name(url: str, title: str = "") -> str:
    try:
        domain = url.split("//")[-1].split("/")[0]
        name = domain.replace("www.", "").split(".")[0]

        known_names = {
            "ibm": "IBM",
            "who": "WHO",
            "nih": "NIH",
            "ncbi": "NCBI",
            "pmc": "PMC",
            "cdc": "CDC",
            "fda": "FDA",
            "microsoft": "Microsoft",
            "google": "Google",
            "openai": "OpenAI",
            "nature": "Nature",
            "sciencedirect": "ScienceDirect",
            "wikipedia": "Wikipedia",
            "arxiv": "arXiv",
            "youtube": "YouTube",
            "youtu": "YouTube",
        }

        return known_names.get(name.lower(), name.upper())

    except Exception:
        return title[:20] if title else "Source"


def is_trusted_url(url: str) -> bool:
    if not url or not url.startswith("http"):
        return True

    lowered_url = url.lower()

    return any(domain in lowered_url for domain in TRUSTED_DOMAINS)


def rate_source_reliability(source: Dict[str, Any]) -> Dict[str, Any]:
    source_type = source.get("type", "unknown")
    url = source.get("url", "")
    domain = get_domain(url)

    if source_type in ["pdf", "docx", "pptx", "image"]:
        return {
            "label": "User-provided source",
            "score": 80,
            "reason": "Uploaded directly by the user; useful but should still be checked for origin and accuracy.",
        }

    if any(domain.endswith(d) or d in domain for d in TRUSTED_DOMAINS):
        return {
            "label": "High trust",
            "score": 95,
            "reason": "Academic, government, medical, or research-oriented source.",
        }

    if any(d in domain for d in MEDIUM_TRUST_DOMAINS):
        return {
            "label": "Moderate trust",
            "score": 75,
            "reason": "Recognized organization or platform, but claims should still be verified.",
        }

    if source_type in ["web_search", "url", "video"]:
        return {
            "label": "Needs review",
            "score": 60,
            "reason": "Web source; reliability depends on publisher, evidence, and author credibility.",
        }

    return {
        "label": "Unknown",
        "score": 50,
        "reason": "Reliability could not be determined automatically.",
    }


def enrich_sources_with_reliability(sources: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    enriched_sources = []

    for source in sources:
        reliability = rate_source_reliability(source)

        enriched_sources.append(
            {
                **source,
                "reliability_label": reliability["label"],
                "reliability_score": reliability["score"],
                "reliability_reason": reliability["reason"],
                "domain": get_domain(source.get("url", "")),
            }
        )

    return enriched_sources


def rank_sources_by_reliability(sources: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    enriched_sources = enrich_sources_with_reliability(sources)

    return sorted(
        enriched_sources,
        key=lambda src: src.get("reliability_score", 0),
        reverse=True,
    )


def compute_overall_trust_score(sources: List[Dict[str, Any]]) -> Dict[str, Any]:
    if not sources:
        return {
            "score": 0,
            "label": "No sources",
            "message": "No sources available for trust scoring.",
        }

    scores = [
        source.get("reliability_score", rate_source_reliability(source)["score"])
        for source in sources
    ]

    avg_score = round(sum(scores) / len(scores))

    if avg_score >= 85:
        label = "Strong"
        message = "The response is supported by generally strong sources."
    elif avg_score >= 70:
        label = "Moderate"
        message = "The response has decent support, but some claims should be verified."
    elif avg_score >= 50:
        label = "Limited"
        message = "The response uses sources that need careful review."
    else:
        label = "Weak"
        message = "The response has weak or insufficient source support."

    return {
        "score": avg_score,
        "label": label,
        "message": message,
    }


def estimate_token_usage(text: str) -> int:
    if not text:
        return 0

    return max(1, round(len(text.split()) * 1.3))


def estimate_workflow_cost(input_tokens: int, output_tokens: int) -> str:
    total_tokens = input_tokens + output_tokens

    if total_tokens < 5000:
        return "Very low"
    if total_tokens < 15000:
        return "Low"
    if total_tokens < 30000:
        return "Medium"

    return "High"


def get_youtube_video_id(url: str) -> str:
    if "youtu.be/" in url:
        return url.split("youtu.be/")[-1].split("?")[0].split("&")[0]

    if "v=" in url:
        return url.split("v=")[-1].split("&")[0]

    match = re.search(r"embed/([^?&/]+)", url)

    if match:
        return match.group(1)

    return ""


@st.cache_data(ttl=3600, show_spinner=False)
def extract_youtube_transcript_cached(url: str, max_chars: int) -> Dict[str, Any]:
    video_id = get_youtube_video_id(url)

    if not video_id:
        return {
            "title": "YouTube Video",
            "source_name": "YouTube",
            "url": url,
            "text": "Could not detect a valid YouTube video ID.",
            "type": "video_error",
        }

    try:
        transcript = YouTubeTranscriptApi.get_transcript(
            video_id,
            languages=["en"],
        )

        text = " ".join([item["text"] for item in transcript])

        if text.strip():
            return {
                "title": f"YouTube Video Transcript ({video_id})",
                "source_name": "YouTube",
                "url": url,
                "text": text[:max_chars],
                "type": "video",
            }

    except Exception:
        pass

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            output_template = os.path.join(temp_dir, "%(id)s.%(ext)s")

            ydl_opts = {
                "skip_download": True,
                "writesubtitles": True,
                "writeautomaticsub": True,
                "subtitleslangs": ["en", "en-US"],
                "subtitlesformat": "vtt",
                "outtmpl": output_template,
                "quiet": True,
                "no_warnings": True,
            }

            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                title = info.get("title", f"YouTube Video ({video_id})")

            transcript_parts = []

            for file_name in os.listdir(temp_dir):
                if file_name.endswith(".vtt"):
                    vtt_path = os.path.join(temp_dir, file_name)

                    for caption in webvtt.read(vtt_path):
                        caption_text = caption.text.replace("\n", " ").strip()

                        if caption_text:
                            transcript_parts.append(caption_text)

            text = " ".join(transcript_parts)

            if text.strip():
                return {
                    "title": title,
                    "source_name": "YouTube",
                    "url": url,
                    "text": text[:max_chars],
                    "type": "video",
                }

    except Exception as e:
        return {
            "title": "YouTube Video",
            "source_name": "YouTube",
            "url": url,
            "text": f"Could not fetch captions/transcript for this YouTube video. Technical detail: {str(e)}",
            "type": "video_error",
        }

    return {
        "title": "YouTube Video",
        "source_name": "YouTube",
        "url": url,
        "text": "No accessible captions or transcript were found for this video.",
        "type": "video_error",
    }


@st.cache_data(ttl=3600, show_spinner=False)
def fetch_url_text_cached(url: str, max_chars: int) -> Dict[str, Any]:
    try:
        response = requests.get(
            url,
            headers={"User-Agent": "Mozilla/5.0"},
            timeout=12,
        )

        response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")

        title = (
            soup.title.string.strip()
            if soup.title and soup.title.string
            else url
        )

        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()

        paragraphs = [
            p.get_text(" ", strip=True)
            for p in soup.find_all("p")
        ]

        text = " ".join(p for p in paragraphs if p)

        if not text:
            text = soup.get_text(" ", strip=True)

        return {
            "title": title,
            "source_name": get_source_name(url, title),
            "url": url,
            "text": text[:max_chars],
            "type": "url",
        }

    except Exception as e:
        return {
            "title": f"Failed to read URL: {url}",
            "source_name": "Source",
            "url": url,
            "text": f"Error: {str(e)}",
            "type": "url_error",
        }


def extract_pdf_text(uploaded_pdf, max_chars: int) -> Dict[str, Any]:
    try:
        pdf_bytes = uploaded_pdf.read()

        doc = fitz.open(
            stream=pdf_bytes,
            filetype="pdf",
        )

        text = "\n".join([page.get_text() for page in doc])
        title = uploaded_pdf.name

        source_name = (
            title.replace(".pdf", "")
            .replace("_", " ")
            .replace("-", " ")
            .strip()[:25]
        ) or "PDF"

        return {
            "title": title,
            "source_name": source_name,
            "url": title,
            "text": text[:max_chars],
            "type": "pdf",
        }

    except Exception as e:
        return {
            "title": uploaded_pdf.name,
            "source_name": "PDF",
            "url": uploaded_pdf.name,
            "text": f"Error reading PDF: {str(e)}",
            "type": "pdf_error",
        }


def extract_docx_text(uploaded_docx, max_chars: int) -> Dict[str, Any]:
    try:
        doc = Document(uploaded_docx)

        paragraphs = [
            para.text
            for para in doc.paragraphs
            if para.text.strip()
        ]

        text = "\n".join(paragraphs)
        title = uploaded_docx.name

        source_name = (
            title.replace(".docx", "")
            .replace("_", " ")
            .replace("-", " ")
            .strip()[:25]
        ) or "DOCX"

        return {
            "title": title,
            "source_name": source_name,
            "url": title,
            "text": text[:max_chars],
            "type": "docx",
        }

    except Exception as e:
        return {
            "title": uploaded_docx.name,
            "source_name": "DOCX",
            "url": uploaded_docx.name,
            "text": f"Error reading DOCX: {str(e)}",
            "type": "docx_error",
        }


def extract_pptx_text(uploaded_pptx, max_chars: int) -> Dict[str, Any]:
    try:
        prs = Presentation(uploaded_pptx)
        slide_texts = []

        for slide_number, slide in enumerate(prs.slides, start=1):
            current_slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    current_slide_text.append(shape.text.strip())

            if current_slide_text:
                slide_texts.append(
                    f"Slide {slide_number}:\n" + "\n".join(current_slide_text)
                )

        text = "\n\n".join(slide_texts)
        title = uploaded_pptx.name

        source_name = (
            title.replace(".pptx", "")
            .replace("_", " ")
            .replace("-", " ")
            .strip()[:25]
        ) or "PPTX"

        return {
            "title": title,
            "source_name": source_name,
            "url": title,
            "text": text[:max_chars] if text else "No readable text found in this PowerPoint file.",
            "type": "pptx",
        }

    except Exception as e:
        return {
            "title": uploaded_pptx.name,
            "source_name": "PPTX",
            "url": uploaded_pptx.name,
            "text": f"Error reading PPTX: {str(e)}",
            "type": "pptx_error",
        }


def extract_image_text(uploaded_image, max_chars: int) -> Dict[str, Any]:
    try:
        image_bytes = uploaded_image.read()
        mime_type = uploaded_image.type or "image/png"

        prompt = """
You are Agent ARCA's image reading tool.

Read and understand this image. Extract:
- visible text
- key information
- table or chart details if present
- important visual observations
- anything useful for research or studying

Return concise but complete extracted content.
""".strip()

        extracted_text = safe_generate_multimodal_content(
            [
                prompt,
                {
                    "mime_type": mime_type,
                    "data": image_bytes,
                },
            ]
        )

        title = uploaded_image.name

        source_name = (
            title.replace(".png", "")
            .replace(".jpg", "")
            .replace(".jpeg", "")
            .replace("_", " ")
            .replace("-", " ")
            .strip()[:25]
        ) or "IMAGE"

        return {
            "title": title,
            "source_name": source_name,
            "url": title,
            "text": extracted_text[:max_chars],
            "type": "image",
        }

    except Exception as e:
        return {
            "title": uploaded_image.name,
            "source_name": "IMAGE",
            "url": uploaded_image.name,
            "text": f"Error reading image: {str(e)}",
            "type": "image_error",
        }


@st.cache_data(ttl=1800, show_spinner=False)
def search_web_cached(
    query: str,
    max_results: int = 3,
    prefer_trusted_sources: bool = False,
    research_depth: str = "Fast Research",
) -> List[Dict[str, Any]]:
    client = get_tavily_client()

    if research_depth == "Deep Research":
        search_depth = "advanced"
        include_raw_content = True
        max_chars = 10000
    else:
        search_depth = "basic"
        include_raw_content = False
        max_chars = 5000

    search_result = client.search(
        query=query,
        search_depth=search_depth,
        max_results=max_results,
        include_answer=False,
        include_raw_content=include_raw_content,
    )

    all_sources = []

    for item in search_result.get("results", []):
        title = item.get("title") or item.get("url", "Untitled Source")
        url = item.get("url", "")
        content = item.get("raw_content") or item.get("content") or ""

        all_sources.append(
            {
                "title": title,
                "source_name": get_source_name(url, title),
                "url": url,
                "text": content[:max_chars],
                "type": "web_search",
            }
        )

    if not prefer_trusted_sources:
        return all_sources

    trusted_sources = [
        src for src in all_sources
        if is_trusted_url(src.get("url", ""))
    ]

    if trusted_sources:
        return trusted_sources

    return all_sources


@st.cache_data(ttl=3600, show_spinner=False)
def find_reference_images_cached(
    query: str,
    max_images: int = 1,
    prefer_trusted_sources: bool = True,
) -> List[Dict[str, str]]:
    try:
        client = get_tavily_client()

        search_result = client.search(
            query=query,
            search_depth="basic",
            max_results=3,
            include_answer=False,
            include_images=True,
            include_image_descriptions=True,
        )

        image_items = []

        top_level_images = search_result.get("images", [])

        for image in top_level_images:
            if isinstance(image, dict):
                image_url = image.get("url", "")
                description = image.get("description", "")
            else:
                image_url = str(image)
                description = ""

            if image_url:
                image_items.append(
                    {
                        "image_url": image_url,
                        "source_url": image_url,
                        "description": description or "Reference image related to the research topic.",
                    }
                )

        for result in search_result.get("results", []):
            result_url = result.get("url", "")
            result_title = result.get("title", "Source")
            result_images = result.get("images", [])

            for image in result_images:
                if isinstance(image, dict):
                    image_url = image.get("url", "")
                    description = image.get("description", "")
                else:
                    image_url = str(image)
                    description = ""

                if not image_url:
                    continue

                if prefer_trusted_sources and result_url.startswith("http"):
                    if not is_trusted_url(result_url):
                        continue

                image_items.append(
                    {
                        "image_url": image_url,
                        "source_url": result_url,
                        "description": description or result_title,
                    }
                )

        unique_images = []
        seen_urls = set()

        for item in image_items:
            image_url = item["image_url"]

            if image_url in seen_urls:
                continue

            seen_urls.add(image_url)
            unique_images.append(item)

            if len(unique_images) >= max_images:
                break

        if not unique_images and prefer_trusted_sources:
            return find_reference_images_cached(
                query=query,
                max_images=max_images,
                prefer_trusted_sources=False,
            )

        return unique_images

    except Exception as e:
        return [
            {
                "image_url": "",
                "source_url": "",
                "description": f"Image search failed: {str(e)[:200]}",
            }
        ]


def build_prompt(
    question: str,
    sources: List[Dict[str, Any]],
    previous_context: str,
    chat_history: List[Any],
    user_mode: str,
    analysis_mode: str,
) -> str:
    source_blocks = []

    for i, src in enumerate(sources, start=1):
        source_blocks.append(
            f"""
[Source {i}]
Source name: {src["source_name"]}
Title: {src["title"]}
URL/File: {src["url"]}
Type: {src["type"]}
Reliability: {src.get("reliability_label", "Unknown")} ({src.get("reliability_score", "N/A")}/100)
Reliability reason: {src.get("reliability_reason", "Not available")}
Content:
{src["text"]}
"""
        )

    sources_text = "\n\n".join(source_blocks)

    recent_chat = chat_history[-6:] if chat_history else []

    chat_text = "\n".join(
        [
            f"{role}: {msg}"
            for role, msg in recent_chat
        ]
    )

    if user_mode == "Student Mode":
        mode_instruction = """
You are operating in Student Mode.

Explain clearly, simply, and help the user learn.

At the end include:
## Quiz
Generate 3 short questions to test understanding.

## Areas to Focus On
Mention concepts the student should revise further.
"""
    else:
        mode_instruction = """
You are operating in General Mode.

Provide a professional research synthesis focused on insights, decision-making, and actionable findings.
"""

    citation_rules = """
Citation rules:
- Use inline markers like [Source 1], [Source 2] for important factual claims.
- Only cite sources that support the claim.
- Do not cite every sentence.
- Prefer higher reliability sources.
- If evidence is weak or missing, say so.
- Do not invent source numbers.
"""

    if analysis_mode == "Compare & Verify":
        output_format = """
Return this format:

## Key Takeaway
One clear sentence with source markers if supported.

## Executive Answer
Clear answer with source markers for important claims.

## Evidence Map
- Claim: ...
  Evidence: [Source X]
  Strength: Strong / Moderate / Weak

## Agreements Across Sources
List agreements with source markers.

## Conflicting Information
List conflicts or say none found.

## Final Consensus
Balanced conclusion.

## Confidence Assessment
Low / Medium / High with a short reason.

## Simple Summary
Very simple explanation.

## Suggested Follow-Up Questions
3 useful questions.

## Recommended Next Steps
3 practical next steps.
"""
    else:
        output_format = """
Return this format:

## Key Takeaway
One clear sentence with source markers if supported.

## Answer
Clear answer with source markers for important factual claims.

## Evidence Map
- Claim: ...
  Evidence: [Source X]
  Strength: Strong / Moderate / Weak

## Key Points
- Point 1
- Point 2
- Point 3

## Simple Summary
Very simple explanation.

## What I verified from sources
Briefly say what was supported.

## Limits
Mention anything unclear or unavailable.

## Suggested Follow-Up Questions
3 useful questions.

## Recommended Next Steps
3 practical next steps.
"""

    return f"""
You are Agent ARCA, a trustworthy AI research assistant.

{mode_instruction}

Analysis mode:
{analysis_mode}

{citation_rules}

User question:
{question}

Previous research context:
{previous_context if previous_context else "No previous context yet."}

Recent conversation:
{chat_text if chat_text else "No previous conversation."}

Sources:
{sources_text}

{output_format}
""".strip()


def evaluate_answer(answer: str, source_count: int) -> Dict[str, float]:
    readability = min(1.0, len(answer) / 1500)
    coverage = min(1.0, source_count / 5)
    confidence = round((readability + coverage) / 2, 2)

    return {
        "readability": round(readability, 2),
        "coverage": round(coverage, 2),
        "confidence": confidence,
    }


def evaluate_quiz_answers(original_answer: str, student_answers: str) -> str:
    prompt = f"""
You are Agent ARCA in Student Mode.

The student studied this material:

{original_answer}

The student answered the quiz like this:

{student_answers}

Evaluate the student's understanding.

Return this format:

## Score
Give a score out of 10.

## What You Got Right
List strengths.

## What Needs Improvement
List weak areas.

## Corrected Explanation
Explain missed concepts simply.

## Suggested Revision Plan
Give 3 focused revision steps.
""".strip()

    return safe_generate_content(prompt)


def clean_export_text(text: str) -> str:
    text = re.sub(r"<[^>]+>", "", text)
    text = text.replace("##", "")
    text = text.replace("*", "")

    return text.strip()


def make_paragraph_safe(text: str) -> str:
    text = clean_export_text(text)
    text = text.replace("&", "&amp;")
    text = text.replace("<", "&lt;")
    text = text.replace(">", "&gt;")

    return text


def extract_bullet_points(answer: str, max_points: int = 8) -> List[str]:
    cleaned = clean_export_text(answer)
    lines = [line.strip() for line in cleaned.split("\n") if line.strip()]

    bullets = []

    skip_starts = [
        "key takeaway",
        "answer",
        "simple summary",
        "evidence map",
        "what i verified",
        "limits",
        "suggested",
        "recommended",
        "quiz",
        "areas to focus",
    ]

    for line in lines:
        lowered = line.lower()

        if any(lowered.startswith(skip) for skip in skip_starts):
            continue

        if len(line) < 4:
            continue

        bullets.append(line)

    return bullets[:max_points] if bullets else ["No findings available."]


def generate_word_report(
    question: str,
    answer: str,
    sources: List[Dict[str, Any]],
    user_mode: str,
    analysis_mode: str,
    quiz_feedback: str = "",
):
    clean_answer = clean_export_text(answer)
    clean_feedback = clean_export_text(quiz_feedback) if quiz_feedback else ""
    trust = compute_overall_trust_score(sources)

    doc = Document()

    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Calibri"
    normal_style.font.size = Pt(11)

    title = doc.add_heading("Agent ARCA Research Report", level=0)
    title.alignment = 1

    subtitle = doc.add_paragraph()
    subtitle.alignment = 1
    subtitle.add_run(
        "AI-powered multimodal research, learning, and decision-intelligence assistant"
    ).italic = True

    doc.add_paragraph("")

    doc.add_heading("Executive Overview", level=1)
    doc.add_paragraph(f"Question: {question}")
    doc.add_paragraph(f"Mode: {user_mode}")
    doc.add_paragraph(f"Analysis Type: {analysis_mode}")
    doc.add_paragraph(f"Sources Used: {len(sources)}")
    doc.add_paragraph(f"Overall Trust Score: {trust['score']}/100 ({trust['label']})")

    doc.add_heading("Research Findings", level=1)
    doc.add_paragraph(clean_answer)

    doc.add_heading("Key Highlights", level=1)

    for point in extract_bullet_points(answer, max_points=6):
        doc.add_paragraph(point, style="List Bullet")

    if clean_feedback:
        doc.add_heading("Student Quiz Feedback", level=1)
        doc.add_paragraph(clean_feedback)

    doc.add_heading("Sources Used", level=1)

    for i, src in enumerate(sources, start=1):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(f"Source {i}: {src['source_name']} — ").bold = True
        p.add_run(src["title"])
        p.add_run(f"\n{src['url']}")
        p.add_run(
            f"\nReliability: {src.get('reliability_label', 'Unknown')} "
            f"({src.get('reliability_score', 'N/A')}/100)"
        )

    doc.add_heading("Recommended Next Steps", level=1)

    for step in [
        "Review the generated findings and summary.",
        "Verify important claims using the listed sources.",
        "Ask follow-up questions in ARCA for deeper analysis.",
        "Export a deck or report for sharing or presentation.",
    ]:
        doc.add_paragraph(step, style="List Number")

    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    return buffer


def generate_pdf_report(
    question: str,
    answer: str,
    sources: List[Dict[str, Any]],
    user_mode: str,
    analysis_mode: str,
    quiz_feedback: str = "",
):
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)

    styles = getSampleStyleSheet()
    story = []
    trust = compute_overall_trust_score(sources)

    clean_answer = make_paragraph_safe(answer)
    clean_feedback = make_paragraph_safe(quiz_feedback) if quiz_feedback else ""

    story.append(Paragraph("Agent ARCA Research Report", styles["Title"]))
    story.append(Spacer(1, 8))

    story.append(
        Paragraph(
            "AI-powered multimodal research, learning, and decision-intelligence assistant",
            styles["BodyText"],
        )
    )

    story.append(Spacer(1, 16))

    story.append(Paragraph("Executive Overview", styles["Heading1"]))
    story.append(Paragraph(f"Question: {make_paragraph_safe(question)}", styles["BodyText"]))
    story.append(Paragraph(f"Mode: {user_mode}", styles["BodyText"]))
    story.append(Paragraph(f"Analysis Type: {analysis_mode}", styles["BodyText"]))
    story.append(Paragraph(f"Sources Used: {len(sources)}", styles["BodyText"]))
    story.append(Paragraph(f"Trust Score: {trust['score']}/100 ({trust['label']})", styles["BodyText"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph("Research Findings", styles["Heading1"]))

    for paragraph in clean_answer.split("\n"):
        if paragraph.strip():
            story.append(Paragraph(paragraph.strip(), styles["BodyText"]))
            story.append(Spacer(1, 6))

    story.append(Spacer(1, 10))
    story.append(Paragraph("Key Highlights", styles["Heading1"]))

    for point in extract_bullet_points(answer, max_points=6):
        story.append(Paragraph(f"• {make_paragraph_safe(point)}", styles["BodyText"]))
        story.append(Spacer(1, 4))

    if clean_feedback:
        story.append(PageBreak())
        story.append(Paragraph("Student Quiz Feedback", styles["Heading1"]))

        for paragraph in clean_feedback.split("\n"):
            if paragraph.strip():
                story.append(Paragraph(paragraph.strip(), styles["BodyText"]))
                story.append(Spacer(1, 6))

    story.append(PageBreak())
    story.append(Paragraph("Sources Used", styles["Heading1"]))

    for i, src in enumerate(sources, start=1):
        source_text = (
            f"Source {i}: {src['source_name']} — {src['title']}<br/>"
            f"{src['url']}<br/>"
            f"Reliability: {src.get('reliability_label', 'Unknown')} "
            f"({src.get('reliability_score', 'N/A')}/100)"
        )

        story.append(Paragraph(make_paragraph_safe(source_text), styles["BodyText"]))
        story.append(Spacer(1, 8))

    story.append(Spacer(1, 8))
    story.append(Paragraph("Recommended Next Steps", styles["Heading1"]))

    for step in [
        "Review the generated findings and summary.",
        "Verify important claims using the listed sources.",
        "Ask follow-up questions in ARCA for deeper analysis.",
        "Export a deck or report for sharing or presentation.",
    ]:
        story.append(Paragraph(f"• {step}", styles["BodyText"]))
        story.append(Spacer(1, 4))

    doc.build(story)
    buffer.seek(0)

    return buffer


def generate_powerpoint_deck(
    question: str,
    answer: str,
    sources: List[Dict[str, Any]],
    user_mode: str,
    analysis_mode: str,
    quiz_feedback: str = "",
):
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    trust = compute_overall_trust_score(sources)
    summary_points = extract_bullet_points(answer, max_points=6)

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Agent ARCA Research Deck"
    slide.placeholders[1].text = f"{user_mode} | {analysis_mode}"

    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Research Question"
    slide.placeholders[1].text = question

    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Executive Summary"

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, point in enumerate(summary_points[:5]):
        paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        paragraph.text = point
        paragraph.level = 0
        paragraph.font.size = PPTPt(18)

    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Research Quality Snapshot"
    slide.placeholders[1].text = (
        f"Sources Used: {len(sources)}\n"
        f"Trust Score: {trust['score']}/100 ({trust['label']})\n"
        f"Mode: {user_mode}\n"
        f"Analysis Type: {analysis_mode}"
    )

    source_lines = []

    for i, src in enumerate(sources, start=1):
        source_lines.append(
            f"Source {i}: {src['source_name']} — "
            f"{src.get('reliability_label', 'Unknown')}"
        )

    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Sources Used"
    slide.placeholders[1].text = "\n".join(source_lines[:10]) if source_lines else "No sources available."

    if quiz_feedback:
        feedback_points = extract_bullet_points(quiz_feedback, max_points=6)
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "Student Quiz Feedback"
        slide.placeholders[1].text = "\n".join(feedback_points)

    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "Recommended Next Steps"
    slide.placeholders[1].text = (
        "1. Review the findings.\n"
        "2. Verify critical claims using the source slide.\n"
        "3. Ask follow-up questions in ARCA.\n"
        "4. Use this deck for study, reporting, or discussion."
    )

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer


def run_pipeline(
    question: str,
    urls: List[str],
    uploaded_pdfs,
    uploaded_docs,
    uploaded_ppts,
    uploaded_images,
    mode: str,
    user_mode: str,
    analysis_mode: str,
    previous_context: str = None,
    chat_history: List[Any] = None,
    max_web_results: int = 3,
    prefer_trusted_sources: bool = False,
    research_depth: str = "Fast Research",
) -> Dict[str, Any]:
    start_time = time.time()
    all_sources = []
    limits = get_limits(research_depth)

    if mode in ["Use my sources", "Use my sources + web search"]:
        for url in urls:
            if "youtube.com" in url or "youtu.be" in url:
                all_sources.append(
                    extract_youtube_transcript_cached(
                        url=url,
                        max_chars=limits["video_chars"],
                    )
                )
            else:
                all_sources.append(
                    fetch_url_text_cached(
                        url=url,
                        max_chars=limits["url_chars"],
                    )
                )

        if uploaded_pdfs:
            for pdf in uploaded_pdfs:
                all_sources.append(
                    extract_pdf_text(
                        uploaded_pdf=pdf,
                        max_chars=limits["file_chars"],
                    )
                )

        if uploaded_docs:
            for docx in uploaded_docs:
                all_sources.append(
                    extract_docx_text(
                        uploaded_docx=docx,
                        max_chars=limits["file_chars"],
                    )
                )

        if uploaded_ppts:
            for pptx in uploaded_ppts:
                all_sources.append(
                    extract_pptx_text(
                        uploaded_pptx=pptx,
                        max_chars=limits["file_chars"],
                    )
                )

        if uploaded_images:
            for image in uploaded_images:
                all_sources.append(
                    extract_image_text(
                        uploaded_image=image,
                        max_chars=limits["image_chars"],
                    )
                )

    should_search_web = (
        mode in ["Search the web", "Use my sources + web search"]
        or len(all_sources) == 0
    )

    if should_search_web:
        all_sources.extend(
            search_web_cached(
                query=question,
                max_results=max_web_results,
                prefer_trusted_sources=prefer_trusted_sources,
                research_depth=research_depth,
            )
        )

    ranked_sources = rank_sources_by_reliability(all_sources)

    prompt = build_prompt(
        question=question,
        sources=ranked_sources,
        previous_context=previous_context or "",
        chat_history=chat_history or [],
        user_mode=user_mode,
        analysis_mode=analysis_mode,
    )

    input_tokens_estimate = estimate_token_usage(prompt)
    answer = safe_generate_content(prompt)
    output_tokens_estimate = estimate_token_usage(answer)

    response_time_seconds = round(time.time() - start_time, 2)
    trust = compute_overall_trust_score(ranked_sources)
    workflow_type = f"{user_mode} | {analysis_mode} | {mode} | {research_depth}"

    new_context = f"""
Last question: {question}

Workflow: {workflow_type}

Analysis mode: {analysis_mode}

Last answer:
{answer[:3000]}

Sources used:
{", ".join([src["source_name"] for src in ranked_sources])}
""".strip()

    return {
        "answer": answer,
        "sources": [
            {
                "title": src["title"],
                "source_name": src["source_name"],
                "url": src["url"],
                "type": src["type"],
                "domain": src.get("domain", ""),
                "reliability_label": src.get("reliability_label", "Unknown"),
                "reliability_score": src.get("reliability_score", 0),
                "reliability_reason": src.get("reliability_reason", ""),
            }
            for src in ranked_sources
        ],
        "context": new_context,
        "evaluation": evaluate_answer(answer, len(ranked_sources)),
        "metadata": {
            "response_time_seconds": response_time_seconds,
            "input_tokens_estimate": input_tokens_estimate,
            "output_tokens_estimate": output_tokens_estimate,
            "total_tokens_estimate": input_tokens_estimate + output_tokens_estimate,
            "estimated_cost_level": estimate_workflow_cost(
                input_tokens_estimate,
                output_tokens_estimate,
            ),
            "workflow_type": workflow_type,
            "research_depth": research_depth,
            "trust_score": trust["score"],
            "trust_label": trust["label"],
            "trust_message": trust["message"],
            "source_count": len(ranked_sources),
        },
    }
